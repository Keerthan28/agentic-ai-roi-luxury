"""
AI Implementation Advisor — generates phased implementation roadmaps
based on user inputs (AI use case, budget, timeline, constraints)
and the ROI model context.  Includes web-sourced OSINT, guardrails,
assumptions disclosure, and full source attribution.
"""

import io
import json
import re
import sys
import traceback

import boto3
import streamlit as st
from botocore.config import Config
from botocore.exceptions import ClientError, NoCredentialsError
from duckduckgo_search import DDGS
from fpdf import FPDF
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BEDROCK_MODEL_ID = "amazon.nova-pro-v1:0"
BEDROCK_REGION = "us-east-1"

SYSTEM_PROMPT = """\
You are a senior AI strategy consultant at a top-tier consulting firm specializing \
in luxury retail digital transformation. You produce implementation roadmaps that \
are so detailed an engineering team could start executing immediately without \
asking clarifying questions.

CRITICAL RULES:
1. SOURCE ATTRIBUTION — Every claim involving a statistic, benchmark, cost estimate, \
   or industry fact MUST cite its source in parentheses, e.g. "(Bain & Company, 2024)" \
   or "(McKinsey Digital, 2024)". If the data comes from the OSINT research snippets \
   provided in the prompt, cite them as shown. If you use your own training knowledge, \
   cite the original publisher and year. NEVER present numbers without attribution.
2. ASSUMPTIONS — You MUST include an "assumptions" array that explicitly lists every \
   assumption you made (market conditions, data availability, team capabilities, \
   technology readiness, etc.). Each assumption must explain what you assumed AND \
   how the roadmap would change if that assumption is wrong.
3. DETAIL LEVEL — Every objective states WHAT, WHO, HOW. Every deliverable is a \
   concrete artifact with quality criteria. Every step is assignable by a PM. \
   Resources specify exact roles/counts. Risks pair threat + mitigation. \
   Success metrics are quantitative with targets.

You MUST return a valid JSON object with this exact structure (no markdown fences, \
no extra text outside JSON):

{
  "executive_summary": "3-4 sentence strategic summary with business impact numbers and source citations",
  "total_estimated_cost": "$X - $Y",
  "expected_roi_timeline": "X-Y months to breakeven",
  "projected_roi": "Y% projected ROI by month Z with brief explanation",
  "phases": [
    {
      "phase_number": 1,
      "phase_name": "Phase Name",
      "duration": "X months",
      "estimated_cost": "$X - $Y",
      "objectives": [
        "Specific objective with source citation where data is referenced (Source, Year)"
      ],
      "steps": [
        {
          "step_number": 1,
          "title": "Short step title",
          "description": "Detailed 2-3 sentence description with tools/platforms, responsible role, and expected output",
          "duration": "X weeks",
          "owner": "Role responsible"
        }
      ],
      "key_deliverables": [
        "Concrete deliverable with quality criteria"
      ],
      "resources_needed": [
        "Role with count and skills (e.g. '2x ML Engineers — Python, PyTorch')"
      ],
      "risks": [
        "Risk: threat description (Source if applicable) | Mitigation: concrete action"
      ],
      "success_metrics": [
        "Quantitative metric with target and benchmark source (e.g. 'Model AUC > 0.85 — industry avg 0.78 per Gartner 2024')"
      ]
    }
  ],
  "technology_stack": ["Technology with version/tier"],
  "quick_wins": [
    "Actionable quick win with timeline, impact, and source for the impact estimate"
  ],
  "critical_success_factors": [
    "Factor with explanation and source if data-backed"
  ],
  "key_recommendations": [
    "Recommendation with rationale and supporting data source"
  ],
  "assumptions": [
    {
      "assumption": "What was assumed",
      "impact_if_wrong": "How the roadmap would change if this assumption is incorrect",
      "confidence": "High/Medium/Low"
    }
  ],
  "sources": [
    {
      "id": "1",
      "reference": "Full source reference (e.g. 'Bain & Company — Luxury Market Report 2024')",
      "url": "URL if available, otherwise empty string",
      "used_for": "Brief description of what data was drawn from this source"
    }
  ]
}

Produce 3-5 phases. Each phase MUST have 4-8 detailed steps. Be specific to the \
luxury goods industry — reference real platforms, real ML techniques, realistic timelines.
"""


# ── Web Search (OSINT) ────────────────────────────────────────────────────

def search_osint(use_case: str, industry: str = "luxury goods") -> list[dict]:
    """Run DuckDuckGo searches to gather recent OSINT for the LLM context."""
    queries = [
        f"{use_case} AI implementation cost benchmark {industry} 2024 2025",
        f"{use_case} ROI statistics {industry} retail",
        f"AI {use_case} luxury retail case study results",
    ]
    results = []
    for q in queries:
        try:
            with DDGS() as ddgs:
                hits = list(ddgs.text(q, max_results=3))
                for h in hits:
                    results.append({
                        "title": h.get("title", ""),
                        "snippet": h.get("body", ""),
                        "url": h.get("href", ""),
                        "query": q,
                    })
        except Exception:
            pass
    seen_urls = set()
    deduped = []
    for r in results:
        if r["url"] not in seen_urls:
            seen_urls.add(r["url"])
            deduped.append(r)
    return deduped[:10]


def format_osint_context(osint_results: list[dict]) -> str:
    """Format search results into a context block for the LLM prompt."""
    if not osint_results:
        return ""
    lines = ["OSINT RESEARCH (cite these sources where you use their data):"]
    for i, r in enumerate(osint_results, 1):
        lines.append(f"[{i}] {r['title']}")
        lines.append(f"    Source: {r['url']}")
        lines.append(f"    Excerpt: {r['snippet'][:300]}")
        lines.append("")
    return "\n".join(lines)


# ── Guardrails ────────────────────────────────────────────────────────────

GUARDRAIL_WARNINGS = []

INPUT_LIMITS = {
    "budget_min": 10_000,
    "budget_max": 50_000_000,
    "timeline_min": 3,
    "timeline_max": 36,
}


def validate_inputs(budget, timeline, team_size, constraints):
    """Check inputs and return list of warning strings."""
    warnings = []
    if budget < 50_000:
        warnings.append(
            f"Budget (${budget:,.0f}) is very low for an enterprise AI initiative. "
            "Results may be optimistic — typical minimum is $50K-$100K."
        )
    if budget > 10_000_000 and team_size in ("Solo / 1 person", "Small (2-5 people)"):
        warnings.append(
            f"Budget (${budget:,.0f}) is large for a {team_size.lower()} team. "
            "Consider whether the team can absorb this spend effectively."
        )
    if timeline < 6 and budget > 1_000_000:
        warnings.append(
            f"A {timeline}-month timeline with ${budget:,.0f} budget is aggressive. "
            "Expect compressed phases and higher execution risk."
        )
    if not constraints:
        warnings.append(
            "No constraints selected. The roadmap assumes no regulatory, "
            "technical, or organizational constraints — review carefully."
        )
    return warnings


def validate_output(roadmap: dict) -> list[str]:
    """Validate the LLM's output and return list of warning strings."""
    warnings = []
    required = ["executive_summary", "total_estimated_cost", "phases"]
    for key in required:
        if key not in roadmap:
            warnings.append(f"Missing required field: {key}")

    phases = roadmap.get("phases", [])
    if len(phases) < 2:
        warnings.append("Fewer than 2 phases generated — roadmap may lack granularity.")
    if len(phases) > 7:
        warnings.append("More than 7 phases — roadmap may be overly fragmented.")

    for ph in phases:
        steps = ph.get("steps", [])
        if len(steps) < 2:
            pn = ph.get("phase_number", "?")
            warnings.append(f"Phase {pn} has fewer than 2 steps — may lack detail.")

    if not roadmap.get("assumptions"):
        warnings.append("No assumptions declared — treat all estimates as unvalidated.")
    if not roadmap.get("sources"):
        warnings.append("No sources cited — data points lack attribution.")

    return warnings

AI_USE_CASES = {
    "Personalized Recommendations": "AI-powered product recommendation engine using customer behavior, preferences, and purchase history to deliver hyper-personalized luxury experiences.",
    "Demand Forecasting": "ML-based demand prediction for inventory optimization, reducing overstock and stockouts across luxury product lines.",
    "Customer Service Automation": "Agentic AI chatbots and virtual concierges for 24/7 luxury customer support with human-like interactions.",
    "Product Authentication": "Computer vision and blockchain-backed AI for authenticating luxury goods and combating counterfeits.",
    "Visual Search & Virtual Try-On": "AR/AI-powered virtual try-on experiences and visual product search for luxury fashion and accessories.",
    "Dynamic Pricing Optimization": "AI-driven pricing strategies that optimize revenue while maintaining brand prestige and exclusivity.",
    "Customer Churn Prediction": "Predictive models to identify at-risk high-value customers and trigger proactive retention campaigns.",
    "Supply Chain Optimization": "AI-powered supply chain visibility, logistics optimization, and sustainability tracking for luxury brands.",
    "Sentiment & Brand Analysis": "NLP-driven monitoring of brand perception, customer sentiment, and competitive intelligence across channels.",
    "Custom / Other": "Describe your specific AI use case in the additional details field.",
}

TEAM_SIZE_OPTIONS = [
    "Solo / 1 person",
    "Small (2-5 people)",
    "Medium (6-15 people)",
    "Large (16-50 people)",
    "Enterprise (50+ people)",
]

CONSTRAINT_OPTIONS = [
    "Data privacy / GDPR compliance",
    "Limited technical team",
    "Legacy system integration required",
    "Must maintain luxury brand experience",
    "Multi-language / multi-region support",
    "Regulatory approval needed",
    "Tight timeline (< 6 months)",
    "Cloud-only (no on-premise)",
    "Must integrate with existing CRM/ERP",
    "Real-time processing required",
]


def get_bedrock_client():
    config = Config(
        region_name=BEDROCK_REGION,
        retries={"max_attempts": 3, "mode": "adaptive"},
    )
    return boto3.client("bedrock-runtime", config=config)


def build_prompt(use_case, use_case_desc, budget, timeline_months, team_size,
                 constraints, additional_details, current_roi_data=None,
                 osint_context=""):
    roi_context = ""
    if current_roi_data:
        roi_context = f"""
Current ROI Model Context (from internal dashboard analytics):
- Current Mean Net ROI: {current_roi_data.get('mean_roi', 'N/A')}x
- AI-Ready Customers (>0.5 readiness): {current_roi_data.get('ai_ready_pct', 'N/A')}
- Average AI Readiness Score: {current_roi_data.get('avg_readiness', 'N/A')}
- Average Revenue Uplift per Customer: ${current_roi_data.get('avg_revenue_uplift', 'N/A'):,.0f}
- Average Retention Savings per Customer: ${current_roi_data.get('avg_retention_savings', 'N/A'):,.0f}
"""

    return f"""Generate an extremely detailed, execution-ready AI implementation roadmap \
for a luxury retail company. The roadmap must be specific enough that a technical \
program manager could begin execution on Day 1 without asking follow-up questions.

PROJECT PARAMETERS:
- Use Case: {use_case}
- Description: {use_case_desc}
- Total Budget: ${budget:,.0f}
- Timeline: {timeline_months} months
- Team Size: {team_size}
- Key Constraints: {', '.join(constraints) if constraints else 'None specified'}
- Additional Details: {additional_details or 'None'}
{roi_context}
{osint_context}

REQUIREMENTS:
1. Produce 3-5 phases that fit within the {timeline_months}-month timeline and ${budget:,.0f} budget.
2. Each phase MUST include 4-8 granular steps with title, description, duration, owner.
3. Deliverables must be concrete artifacts with quality bars.
4. Risks must pair a specific threat with a concrete mitigation action.
5. Success metrics must be quantitative with numeric targets.
6. Reference real-world luxury industry platforms, ML techniques, and integration points.
7. Allocate budget realistically — earlier phases 30-40%, later phases 15-20%.
8. SOURCE ATTRIBUTION IS MANDATORY — cite every statistic. Use the OSINT research \
   snippets provided above, the ROI model context, and your training knowledge. \
   Include a "sources" array with full references and URLs.
9. ASSUMPTIONS — include an "assumptions" array listing every assumption with \
   confidence level and impact if wrong.

Return ONLY the JSON object — no markdown, no commentary."""


def invoke_bedrock(prompt):
    """Call the LLM and return the parsed roadmap."""
    client = get_bedrock_client()
    response = client.converse(
        modelId=BEDROCK_MODEL_ID,
        messages=[{"role": "user", "content": [{"text": prompt}]}],
        system=[{"text": SYSTEM_PROMPT}],
        inferenceConfig={"maxTokens": 8192, "temperature": 0.4},
    )
    raw = response["output"]["message"]["content"][0]["text"]
    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1]
        if raw.endswith("```"):
            raw = raw[:-3]
    return json.loads(raw)


# ── Export: PDF ────────────────────────────────────────────────────────────

def _safe(text):
    """Sanitize text for FPDF (latin-1 safe)."""
    return str(text).encode("latin-1", "replace").decode("latin-1")


def _bullet(pdf, text, indent=8):
    """Render a bullet item with safe indentation."""
    left = pdf.l_margin
    pdf.set_x(left + indent)
    effective_w = pdf.w - pdf.r_margin - pdf.x
    if effective_w < 10:
        effective_w = pdf.w - pdf.l_margin - pdf.r_margin
        pdf.set_x(left)
    pdf.multi_cell(effective_w, 5, _safe(f"- {text}"))


def generate_pdf(roadmap):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()
    pw = pdf.w - pdf.l_margin - pdf.r_margin

    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(44, 62, 80)
    pdf.cell(pw, 12, _safe("AI Implementation Roadmap"), ln=True, align="C")
    pdf.ln(4)

    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(80, 80, 80)
    pdf.multi_cell(pw, 5, _safe(roadmap.get("executive_summary", "")))
    pdf.ln(4)

    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(44, 62, 80)
    third = pw / 3
    pdf.cell(third, 7, _safe(f"Total Cost: {roadmap.get('total_estimated_cost', 'N/A')}"))
    pdf.cell(third, 7, _safe(f"ROI Timeline: {roadmap.get('expected_roi_timeline', 'N/A')}"))
    projected_roi = roadmap.get("projected_roi", "")
    if projected_roi:
        pdf.cell(third, 7, _safe(f"Projected ROI: {projected_roi}"), ln=True)
    else:
        pdf.cell(third, 7, "", ln=True)
    pdf.ln(4)

    if roadmap.get("quick_wins"):
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(pw, 8, _safe("Quick Wins"), ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(60, 60, 60)
        for win in roadmap["quick_wins"]:
            _bullet(pdf, win)
        pdf.ln(3)

    for phase in roadmap.get("phases", []):
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(41, 128, 185)
        pnum = phase.get("phase_number", "")
        pname = phase.get("phase_name", "")
        pdur = phase.get("duration", "")
        pcost = phase.get("estimated_cost", "")
        pdf.cell(pw, 8, _safe(f"Phase {pnum}: {pname} ({pdur} | {pcost})"), ln=True)

        pdf.set_text_color(60, 60, 60)
        sections = [
            ("Objectives", "objectives"),
            ("Key Deliverables", "key_deliverables"),
            ("Key Activities", "key_activities"),
            ("Resources Needed", "resources_needed"),
            ("Risks & Mitigations", "risks"),
            ("Success Metrics", "success_metrics"),
        ]
        for section, key in sections:
            items = phase.get(key, [])
            if not items:
                continue
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(pw, 6, _safe(section), ln=True)
            pdf.set_font("Helvetica", "", 9)
            for item in items:
                _bullet(pdf, item)

        steps = phase.get("steps", [])
        if steps:
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 11)
            pdf.set_text_color(41, 128, 185)
            pdf.cell(pw, 7, _safe("Implementation Steps"), ln=True)
            pdf.set_text_color(60, 60, 60)
            for step in steps:
                pdf.set_font("Helvetica", "B", 9)
                snum = step.get("step_number", "")
                stitle = step.get("title", "")
                pdf.cell(pw, 5, _safe(f"Step {snum}: {stitle}"), ln=True)
                pdf.set_font("Helvetica", "", 8)
                sdesc = step.get("description", "")
                if sdesc:
                    pdf.set_x(pdf.l_margin + 8)
                    ew = pw - 8
                    pdf.multi_cell(ew, 4, _safe(sdesc))
                sdur = step.get("duration", "")
                sowner = step.get("owner", "")
                if sdur or sowner:
                    pdf.set_font("Helvetica", "I", 8)
                    pdf.set_x(pdf.l_margin + 8)
                    pdf.cell(pw - 8, 4, _safe(f"Duration: {sdur}  |  Owner: {sowner}"), ln=True)
                pdf.ln(1)
        pdf.ln(3)

    if roadmap.get("technology_stack"):
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(pw, 8, _safe("Technology Stack"), ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(60, 60, 60)
        pdf.multi_cell(pw, 5, _safe(", ".join(roadmap["technology_stack"])))
        pdf.ln(3)

    if roadmap.get("critical_success_factors"):
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(pw, 8, _safe("Critical Success Factors"), ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(60, 60, 60)
        for factor in roadmap["critical_success_factors"]:
            _bullet(pdf, factor)

    if roadmap.get("key_recommendations"):
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(pw, 8, _safe("Key Recommendations"), ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(60, 60, 60)
        for rec in roadmap["key_recommendations"]:
            _bullet(pdf, rec)

    assumptions = roadmap.get("assumptions", [])
    if assumptions:
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(pw, 8, _safe("Assumptions"), ln=True)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(60, 60, 60)
        for a in assumptions:
            if isinstance(a, dict):
                conf = a.get("confidence", "Medium")
                txt = f"[{conf}] {a.get('assumption', '')} — If wrong: {a.get('impact_if_wrong', 'N/A')}"
                _bullet(pdf, txt)
            else:
                _bullet(pdf, str(a))

    sources = roadmap.get("sources", [])
    if sources:
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(pw, 8, _safe("Sources & Attribution"), ln=True)
        pdf.set_font("Helvetica", "", 8)
        pdf.set_text_color(60, 60, 60)
        for s in sources:
            if isinstance(s, dict):
                ref = s.get("reference", "")
                url = s.get("url", "")
                used = s.get("used_for", "")
                line = f"{ref}"
                if url:
                    line += f" — {url}"
                if used:
                    line += f" (Used for: {used})"
                _bullet(pdf, line)
            else:
                _bullet(pdf, str(s))

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Export: DOCX ───────────────────────────────────────────────────────────

def generate_docx(roadmap):
    doc = Document()

    style = doc.styles["Heading 1"]
    style.font.color.rgb = RGBColor(44, 62, 80)

    doc.add_heading("AI Implementation Roadmap", level=0)
    doc.add_paragraph(roadmap.get("executive_summary", ""))

    t = doc.add_table(rows=1, cols=3, style="Light Shading Accent 1")
    t.rows[0].cells[0].text = f"Total Cost: {roadmap.get('total_estimated_cost', 'N/A')}"
    t.rows[0].cells[1].text = f"ROI Timeline: {roadmap.get('expected_roi_timeline', 'N/A')}"
    projected_roi = roadmap.get("projected_roi", "")
    t.rows[0].cells[2].text = f"Projected ROI: {projected_roi}" if projected_roi else ""

    if roadmap.get("quick_wins"):
        doc.add_heading("Quick Wins", level=1)
        for win in roadmap["quick_wins"]:
            doc.add_paragraph(win, style="List Bullet")

    for phase in roadmap.get("phases", []):
        pnum = phase.get("phase_number", "")
        pname = phase.get("phase_name", "")
        pdur = phase.get("duration", "")
        pcost = phase.get("estimated_cost", "")
        doc.add_heading(f"Phase {pnum}: {pname} ({pdur} | {pcost})", level=1)
        for section, key in [("Objectives", "objectives"),
                             ("Key Deliverables", "key_deliverables"),
                             ("Key Activities", "key_activities"),
                             ("Resources Needed", "resources_needed"),
                             ("Risks & Mitigations", "risks"),
                             ("Success Metrics", "success_metrics")]:
            items = phase.get(key, [])
            if not items:
                continue
            doc.add_heading(section, level=2)
            for item in items:
                doc.add_paragraph(item, style="List Bullet")

        steps = phase.get("steps", [])
        if steps:
            doc.add_heading("Implementation Steps", level=2)
            for step in steps:
                snum = step.get("step_number", "")
                stitle = step.get("title", "")
                sdesc = step.get("description", "")
                sdur = step.get("duration", "")
                sowner = step.get("owner", "")
                p = doc.add_paragraph()
                run_title = p.add_run(f"Step {snum}: {stitle}")
                run_title.bold = True
                run_title.font.size = Pt(11)
                if sdesc:
                    doc.add_paragraph(sdesc)
                if sdur or sowner:
                    meta = doc.add_paragraph()
                    meta_run = meta.add_run(f"Duration: {sdur}  |  Owner: {sowner}")
                    meta_run.italic = True
                    meta_run.font.size = Pt(9)
                    meta_run.font.color.rgb = RGBColor(0x75, 0x78, 0x7B)

    if roadmap.get("technology_stack"):
        doc.add_heading("Technology Stack", level=1)
        doc.add_paragraph(", ".join(roadmap["technology_stack"]))

    if roadmap.get("critical_success_factors"):
        doc.add_heading("Critical Success Factors", level=1)
        for factor in roadmap["critical_success_factors"]:
            doc.add_paragraph(factor, style="List Bullet")

    if roadmap.get("key_recommendations"):
        doc.add_heading("Key Recommendations", level=1)
        for rec in roadmap["key_recommendations"]:
            doc.add_paragraph(rec, style="List Bullet")

    assumptions = roadmap.get("assumptions", [])
    if assumptions:
        doc.add_heading("Assumptions", level=1)
        for a in assumptions:
            if isinstance(a, dict):
                p = doc.add_paragraph()
                run = p.add_run(f"[{a.get('confidence', 'Medium')}] {a.get('assumption', '')}")
                run.bold = True
                run.font.size = Pt(10)
                impact = a.get("impact_if_wrong", "")
                if impact:
                    doc.add_paragraph(f"If wrong: {impact}")
            else:
                doc.add_paragraph(str(a), style="List Bullet")

    sources = roadmap.get("sources", [])
    if sources:
        doc.add_heading("Sources & Attribution", level=1)
        for s in sources:
            if isinstance(s, dict):
                ref = s.get("reference", "")
                url = s.get("url", "")
                used = s.get("used_for", "")
                p = doc.add_paragraph()
                r_run = p.add_run(ref)
                r_run.bold = True
                r_run.font.size = Pt(9)
                if url:
                    p.add_run(f"\n{url}").font.size = Pt(8)
                if used:
                    p.add_run(f"\nUsed for: {used}").font.size = Pt(8)
            else:
                doc.add_paragraph(str(s), style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Export: PPTX  (Deloitte-themed) ───────────────────────────────────────

# Deloitte brand palette
_DT_GREEN  = PptxRGB(0x86, 0xBC, 0x25)
_DT_BLACK  = PptxRGB(0x00, 0x00, 0x00)
_DT_WHITE  = PptxRGB(0xFF, 0xFF, 0xFF)
_DT_DGRAY  = PptxRGB(0x2D, 0x2D, 0x2D)
_DT_MGRAY  = PptxRGB(0x75, 0x78, 0x7B)
_DT_LGRAY  = PptxRGB(0xF2, 0xF2, 0xF2)
_DT_TEAL   = PptxRGB(0x00, 0xA3, 0xE0)
_DT_NAVY   = PptxRGB(0x01, 0x2B, 0x5D)
_DT_AMBER  = PptxRGB(0xED, 0x8B, 0x00)
_DT_RED    = PptxRGB(0xDA, 0x29, 0x1C)

_PHASE_COLORS = [_DT_GREEN, _DT_TEAL, _DT_NAVY, _DT_AMBER,
                 PptxRGB(0x6F, 0xC2, 0xB0), PptxRGB(0xA0, 0xD5, 0x6C)]

SW = PptxInches(13.333)
SH = PptxInches(7.5)


def _set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_shape_text(shape, text, size=12, bold=False, color=_DT_WHITE,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = PptxPt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass


def _add_green_bar(slide, top=Emu(0)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, SW, PptxInches(0.08))
    _set_shape_fill(bar, _DT_GREEN)
    bar.line.fill.background()
    return bar


def _add_footer(slide, text="Deloitte AI Advisory  |  Confidential"):
    ft = slide.shapes.add_textbox(
        PptxInches(0.5), SH - PptxInches(0.45),
        SW - PptxInches(1), PptxInches(0.35))
    p = ft.text_frame.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(8)
    p.font.color.rgb = _DT_MGRAY
    p.alignment = PP_ALIGN.CENTER


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Black header band with green accent line."""
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SW, PptxInches(1.3))
    _set_shape_fill(hdr, _DT_BLACK)
    hdr.line.fill.background()

    _add_green_bar(slide, top=PptxInches(1.3))

    tb = slide.shapes.add_textbox(
        PptxInches(0.7), PptxInches(0.18), SW - PptxInches(1.4), PptxInches(0.65))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = PptxPt(28)
    p.font.bold = True
    p.font.color.rgb = _DT_WHITE

    if subtitle_text:
        stb = slide.shapes.add_textbox(
            PptxInches(0.7), PptxInches(0.82), SW - PptxInches(1.4), PptxInches(0.4))
        sp = stb.text_frame.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = PptxPt(14)
        sp.font.color.rgb = _DT_GREEN


def _add_bullet_box(slide, left, top, width, height, title, items,
                    title_color=_DT_GREEN, bg=_DT_LGRAY):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _set_shape_fill(box, bg)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(13)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = PptxPt(4)
    for item in items[:8]:
        bp = tf.add_paragraph()
        bp.text = f"  {item}"
        bp.font.size = PptxPt(10)
        bp.font.color.rgb = _DT_DGRAY
        bp.space_before = PptxPt(2)


def _parse_cost_number(cost_str):
    """Extract a numeric dollar value from a string like '$500K - $1M'."""
    nums = re.findall(r'\$?([\d,.]+)\s*([KkMmBb])?', str(cost_str))
    if not nums:
        return 0
    val_str, suffix = nums[0]
    val = float(val_str.replace(',', ''))
    mult = {'k': 1_000, 'm': 1_000_000, 'b': 1_000_000_000}
    val *= mult.get(suffix.lower(), 1) if suffix else 1
    return val


def generate_pptx(roadmap):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    phases = roadmap.get("phases", [])

    # ── SLIDE 1: Title ─────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SW, SH)
    _set_shape_fill(bg, _DT_BLACK)
    bg.line.fill.background()

    _add_green_bar(s1, top=PptxInches(2.6))

    tb = s1.shapes.add_textbox(
        PptxInches(0.9), PptxInches(1.2), PptxInches(10), PptxInches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "AI Implementation Roadmap"
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = _DT_WHITE

    stb = s1.shapes.add_textbox(
        PptxInches(0.9), PptxInches(2.8), PptxInches(11), PptxInches(0.5))
    sp = stb.text_frame.paragraphs[0]
    sp.text = "Luxury Goods Industry  |  Agentic AI Strategy"
    sp.font.size = PptxPt(18)
    sp.font.color.rgb = _DT_GREEN

    desc = s1.shapes.add_textbox(
        PptxInches(0.9), PptxInches(3.8), PptxInches(11), PptxInches(1.5))
    dp = desc.text_frame.paragraphs[0]
    dp.text = roadmap.get("executive_summary", "")
    dp.font.size = PptxPt(13)
    dp.font.color.rgb = _DT_MGRAY
    desc.text_frame.word_wrap = True

    _add_footer(s1, "Deloitte AI Advisory  |  Confidential")

    # ── SLIDE 2: Executive Overview ────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _add_title_bar(s2, "Executive Overview", "Investment & Return Summary")
    _add_footer(s2)

    projected_roi = roadmap.get("projected_roi", "")
    metrics = [
        ("Total Investment", roadmap.get("total_estimated_cost", "N/A"), _DT_GREEN),
        ("ROI Timeline", roadmap.get("expected_roi_timeline", "N/A"), _DT_TEAL),
        ("Projected ROI", projected_roi, _DT_AMBER) if projected_roi
        else ("Phases", str(len(phases)), _DT_NAVY),
    ]
    n_metrics = len(metrics)
    mw = PptxInches(3.5)
    gap = (SW - PptxInches(1.4) - mw * n_metrics) / max(n_metrics - 1, 1)
    for i, (label, value, clr) in enumerate(metrics):
        left = int(PptxInches(0.7)) + i * (int(mw) + int(gap))
        card = s2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, PptxInches(1.8), mw, PptxInches(1.5))
        _set_shape_fill(card, _DT_LGRAY)
        card.line.color.rgb = clr
        card.line.width = PptxPt(2)

        vtb = s2.shapes.add_textbox(left + PptxInches(0.2), PptxInches(2.0),
                                    mw - PptxInches(0.4), PptxInches(0.35))
        vp = vtb.text_frame.paragraphs[0]
        vp.text = label.upper()
        vp.font.size = PptxPt(10)
        vp.font.bold = True
        vp.font.color.rgb = _DT_MGRAY

        nvtb = s2.shapes.add_textbox(left + PptxInches(0.2), PptxInches(2.3),
                                     mw - PptxInches(0.4), PptxInches(0.85))
        nvtb.text_frame.word_wrap = True
        nvp = nvtb.text_frame.paragraphs[0]
        nvp.text = str(value)
        font_sz = PptxPt(22) if len(str(value)) < 20 else PptxPt(16) if len(str(value)) < 35 else PptxPt(13)
        nvp.font.size = font_sz
        nvp.font.bold = True
        nvp.font.color.rgb = clr

    if roadmap.get("quick_wins"):
        qw_title = s2.shapes.add_textbox(
            PptxInches(0.9), PptxInches(3.6), PptxInches(3), PptxInches(0.4))
        qp = qw_title.text_frame.paragraphs[0]
        qp.text = "QUICK WINS"
        qp.font.size = PptxPt(12)
        qp.font.bold = True
        qp.font.color.rgb = _DT_GREEN

        for j, win in enumerate(roadmap["quick_wins"][:4]):
            wtb = s2.shapes.add_textbox(
                PptxInches(1.2), PptxInches(4.05 + j * 0.35),
                PptxInches(11), PptxInches(0.35))
            wp = wtb.text_frame.paragraphs[0]
            wp.text = f"\u2714  {win}"
            wp.font.size = PptxPt(11)
            wp.font.color.rgb = _DT_DGRAY

    # ── SLIDE 3: Phase Timeline Flowchart ──────────────────────────
    if phases:
        s3 = prs.slides.add_slide(blank)
        _add_title_bar(s3, "Implementation Timeline", "Phased Delivery Approach")
        _add_footer(s3)

        n = len(phases)
        avail_w = PptxInches(11.5)
        gap = PptxInches(0.15)
        box_w_emu = int((avail_w - gap * (n - 1)) / n) if n > 0 else avail_w
        box_h = PptxInches(3.8)
        y_top = PptxInches(1.8)
        start_x = PptxInches(0.9)

        for i, phase in enumerate(phases):
            clr = _PHASE_COLORS[i % len(_PHASE_COLORS)]
            left = int(start_x) + i * (box_w_emu + int(gap))

            # Phase column
            col = s3.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, y_top, box_w_emu, box_h)
            _set_shape_fill(col, _DT_LGRAY)
            col.line.color.rgb = clr
            col.line.width = PptxPt(2)

            # Phase number badge
            badge_sz = PptxInches(0.55)
            badge = s3.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left + int(box_w_emu / 2) - int(badge_sz / 2),
                y_top - int(badge_sz / 2),
                badge_sz, badge_sz)
            _set_shape_fill(badge, clr)
            badge.line.fill.background()
            _set_shape_text(badge, str(phase.get("phase_number", i + 1)),
                           size=16, bold=True, color=_DT_WHITE,
                           alignment=PP_ALIGN.CENTER)

            # Phase name
            ntb = s3.shapes.add_textbox(
                left + PptxInches(0.1), y_top + PptxInches(0.4),
                box_w_emu - PptxInches(0.2), PptxInches(0.5))
            ntb.text_frame.word_wrap = True
            np_ = ntb.text_frame.paragraphs[0]
            np_.text = phase.get("phase_name", "")
            np_.font.size = PptxPt(11)
            np_.font.bold = True
            np_.font.color.rgb = _DT_BLACK
            np_.alignment = PP_ALIGN.CENTER

            # Duration & cost
            dtb = s3.shapes.add_textbox(
                left + PptxInches(0.1), y_top + PptxInches(0.85),
                box_w_emu - PptxInches(0.2), PptxInches(0.6))
            dtb.text_frame.word_wrap = True
            ddp = dtb.text_frame.paragraphs[0]
            ddp.text = phase.get("duration", "")
            ddp.font.size = PptxPt(9)
            ddp.font.color.rgb = clr
            ddp.font.bold = True
            ddp.alignment = PP_ALIGN.CENTER
            cdp = dtb.text_frame.add_paragraph()
            cdp.text = phase.get("estimated_cost", "")
            cdp.font.size = PptxPt(9)
            cdp.font.color.rgb = _DT_MGRAY
            cdp.alignment = PP_ALIGN.CENTER

            # Key deliverables
            deliverables = phase.get("key_deliverables") or phase.get("key_activities", [])
            itb = s3.shapes.add_textbox(
                left + PptxInches(0.1), y_top + PptxInches(1.45),
                box_w_emu - PptxInches(0.2), box_h - PptxInches(1.6))
            itb.text_frame.word_wrap = True
            for k, d in enumerate(deliverables[:5]):
                ip = itb.text_frame.paragraphs[0] if k == 0 else itb.text_frame.add_paragraph()
                ip.text = f"\u2022 {d}"
                ip.font.size = PptxPt(8)
                ip.font.color.rgb = _DT_DGRAY
                ip.space_before = PptxPt(2)

            # Arrow connector between phases
            if i < n - 1:
                arr_left = left + box_w_emu
                arr_top = y_top + int(box_h / 2) - PptxInches(0.15)
                arrow = s3.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arr_left, arr_top, int(gap), PptxInches(0.3))
                _set_shape_fill(arrow, _DT_GREEN)
                arrow.line.fill.background()

    # ── SLIDE 4: Budget Breakdown ──────────────────────────────────
    if phases:
        s4 = prs.slides.add_slide(blank)
        _add_title_bar(s4, "Budget Allocation", "Investment Distribution by Phase")
        _add_footer(s4)

        costs = []
        labels = []
        for ph in phases:
            c = _parse_cost_number(ph.get("estimated_cost", "0"))
            costs.append(c if c > 0 else 1)
            labels.append(f"Phase {ph.get('phase_number', '?')}: {ph.get('phase_name', '')}")

        total_c = sum(costs) or 1
        bar_max_w = PptxInches(8.0)
        bar_h = PptxInches(0.6)
        x_start = PptxInches(3.8)
        y_start = PptxInches(1.9)

        for i, (label, cost) in enumerate(zip(labels, costs)):
            y = y_start + i * (bar_h + PptxInches(0.25))
            clr = _PHASE_COLORS[i % len(_PHASE_COLORS)]

            # Label
            ltb = s4.shapes.add_textbox(
                PptxInches(0.7), y, PptxInches(3.0), bar_h)
            lp = ltb.text_frame.paragraphs[0]
            lp.text = label
            lp.font.size = PptxPt(11)
            lp.font.bold = True
            lp.font.color.rgb = _DT_DGRAY
            lp.alignment = PP_ALIGN.RIGHT

            # Bar
            bw = int(bar_max_w * (cost / total_c))
            bw = max(bw, PptxInches(0.5))
            bar = s4.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x_start, y, bw, bar_h)
            _set_shape_fill(bar, clr)
            bar.line.fill.background()

            pct = cost / total_c * 100
            vtb = s4.shapes.add_textbox(
                x_start + PptxInches(0.15), y, bw - PptxInches(0.1), bar_h)
            vp = vtb.text_frame.paragraphs[0]
            cost_txt = phases[i].get("estimated_cost", "")
            vp.text = f"{cost_txt}  ({pct:.0f}%)"
            vp.font.size = PptxPt(10)
            vp.font.bold = True
            vp.font.color.rgb = _DT_WHITE
            try:
                vtb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

    # ── SLIDE 5+: Phase Detail Slides ──────────────────────────────
    for i, phase in enumerate(phases):
        sd = prs.slides.add_slide(blank)
        pnum = phase.get("phase_number", i + 1)
        pname = phase.get("phase_name", "")
        clr = _PHASE_COLORS[i % len(_PHASE_COLORS)]
        _add_title_bar(sd,
                       f"Phase {pnum}: {pname}",
                       f"{phase.get('duration', '')}  |  {phase.get('estimated_cost', '')}")
        _add_footer(sd)

        col_w = PptxInches(5.8)
        col_h = PptxInches(2.3)

        deliverables = phase.get("key_deliverables") or phase.get("key_activities", [])
        _add_bullet_box(sd, PptxInches(0.7), PptxInches(1.7),
                        col_w, col_h, "Objectives",
                        phase.get("objectives", []), title_color=clr)
        _add_bullet_box(sd, PptxInches(6.8), PptxInches(1.7),
                        col_w, col_h, "Key Deliverables",
                        deliverables, title_color=clr)
        _add_bullet_box(sd, PptxInches(0.7), PptxInches(4.2),
                        PptxInches(3.7), col_h, "Resources",
                        phase.get("resources_needed", []), title_color=_DT_TEAL)
        _add_bullet_box(sd, PptxInches(4.7), PptxInches(4.2),
                        PptxInches(3.7), col_h, "Risks",
                        phase.get("risks", []),
                        title_color=_DT_RED, bg=PptxRGB(0xFD, 0xF0, 0xEF))
        _add_bullet_box(sd, PptxInches(8.7), PptxInches(4.2),
                        PptxInches(3.9), col_h, "Success Metrics",
                        phase.get("success_metrics", []), title_color=_DT_GREEN)

        # Steps slide(s) for this phase
        steps = phase.get("steps", [])
        if steps:
            page_size = 4
            for pg_start in range(0, len(steps), page_size):
                page_steps = steps[pg_start:pg_start + page_size]
                ss = prs.slides.add_slide(blank)
                pg_label = f" ({pg_start // page_size + 1})" if len(steps) > page_size else ""
                _add_title_bar(ss,
                               f"Phase {pnum}: Implementation Steps{pg_label}",
                               pname)
                _add_footer(ss)

                card_w = PptxInches(5.9)
                card_h = PptxInches(1.25)
                for si, step in enumerate(page_steps):
                    row, col_idx = divmod(si, 2)
                    left = PptxInches(0.7) + col_idx * (card_w + PptxInches(0.3))
                    top = PptxInches(1.7) + row * (card_h + PptxInches(0.2))

                    card = ss.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
                    _set_shape_fill(card, _DT_LGRAY)
                    card.line.color.rgb = clr
                    card.line.width = PptxPt(1.5)

                    # Step number badge
                    badge_sz = PptxInches(0.4)
                    badge = ss.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        left + PptxInches(0.15),
                        top + PptxInches(0.12),
                        badge_sz, badge_sz)
                    _set_shape_fill(badge, clr)
                    badge.line.fill.background()
                    _set_shape_text(badge, str(step.get("step_number", si + 1)),
                                   size=12, bold=True, color=_DT_WHITE,
                                   alignment=PP_ALIGN.CENTER)

                    # Title
                    ttb = ss.shapes.add_textbox(
                        left + PptxInches(0.65), top + PptxInches(0.08),
                        card_w - PptxInches(0.8), PptxInches(0.3))
                    tp = ttb.text_frame.paragraphs[0]
                    tp.text = step.get("title", "")
                    tp.font.size = PptxPt(11)
                    tp.font.bold = True
                    tp.font.color.rgb = _DT_BLACK

                    # Description
                    dtb = ss.shapes.add_textbox(
                        left + PptxInches(0.65), top + PptxInches(0.38),
                        card_w - PptxInches(0.8), PptxInches(0.55))
                    dtb.text_frame.word_wrap = True
                    ddp = dtb.text_frame.paragraphs[0]
                    ddp.text = step.get("description", "")[:200]
                    ddp.font.size = PptxPt(8)
                    ddp.font.color.rgb = _DT_DGRAY

                    # Duration + Owner
                    mtb = ss.shapes.add_textbox(
                        left + PptxInches(0.65), top + card_h - PptxInches(0.25),
                        card_w - PptxInches(0.8), PptxInches(0.2))
                    mp = mtb.text_frame.paragraphs[0]
                    sdur = step.get("duration", "")
                    sowner = step.get("owner", "")
                    mp.text = f"{sdur}  |  {sowner}"
                    mp.font.size = PptxPt(7)
                    mp.font.italic = True
                    mp.font.color.rgb = _DT_MGRAY

    # ── SLIDE: Technology Stack ────────────────────────────────────
    tech = roadmap.get("technology_stack", [])
    if tech:
        st_slide = prs.slides.add_slide(blank)
        _add_title_bar(st_slide, "Technology Stack", "Recommended Platform & Tools")
        _add_footer(st_slide)

        cols = min(len(tech), 4)
        rows = (len(tech) + cols - 1) // cols
        card_w = PptxInches(2.8)
        card_h = PptxInches(1.0)
        x_off = PptxInches(0.9)
        y_off = PptxInches(2.0)

        for j, t in enumerate(tech):
            r, c_ = divmod(j, cols)
            left = x_off + c_ * (card_w + PptxInches(0.3))
            top = y_off + r * (card_h + PptxInches(0.25))
            card = st_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            _set_shape_fill(card, _DT_LGRAY)
            card.line.color.rgb = _DT_GREEN
            card.line.width = PptxPt(1.5)
            _set_shape_text(card, t, size=11, bold=True,
                           color=_DT_DGRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE: Critical Success Factors / Recommendations ─────────
    csf = roadmap.get("critical_success_factors", [])
    recs = roadmap.get("key_recommendations", [])
    if csf or recs:
        sf = prs.slides.add_slide(blank)
        _add_title_bar(sf, "Critical Success Factors & Recommendations")
        _add_footer(sf)

        if csf:
            _add_bullet_box(sf, PptxInches(0.7), PptxInches(1.7),
                            PptxInches(5.8), PptxInches(4.5),
                            "Critical Success Factors", csf,
                            title_color=_DT_GREEN)
        if recs:
            _add_bullet_box(sf, PptxInches(6.8), PptxInches(1.7),
                            PptxInches(5.8), PptxInches(4.5),
                            "Key Recommendations", recs,
                            title_color=_DT_TEAL)

    # ── SLIDE: Assumptions ─────────────────────────────────────────
    assumptions = roadmap.get("assumptions", [])
    if assumptions:
        sa = prs.slides.add_slide(blank)
        _add_title_bar(sa, "Assumptions & Caveats",
                       "Review and validate before execution")
        _add_footer(sa)
        a_items = []
        for a in assumptions:
            if isinstance(a, dict):
                conf = a.get("confidence", "Medium")
                a_items.append(
                    f"[{conf}] {a.get('assumption', '')} — "
                    f"If wrong: {a.get('impact_if_wrong', 'N/A')}"
                )
            else:
                a_items.append(str(a))
        _add_bullet_box(sa, PptxInches(0.7), PptxInches(1.7),
                        PptxInches(11.9), PptxInches(4.8),
                        "Assumptions", a_items,
                        title_color=_DT_AMBER)

    # ── SLIDE: Sources ───────────────────────────────────────────
    sources = roadmap.get("sources", [])
    if sources:
        ss_slide = prs.slides.add_slide(blank)
        _add_title_bar(ss_slide, "Sources & Attribution",
                       "All data points are attributed to the following sources")
        _add_footer(ss_slide)
        s_items = []
        for s in sources:
            if isinstance(s, dict):
                ref = s.get("reference", "")
                url = s.get("url", "")
                used = s.get("used_for", "")
                line = ref
                if url:
                    line += f"  ({url})"
                if used:
                    line += f"  — {used}"
                s_items.append(line)
            else:
                s_items.append(str(s))
        _add_bullet_box(ss_slide, PptxInches(0.7), PptxInches(1.7),
                        PptxInches(11.9), PptxInches(4.8),
                        "Sources", s_items,
                        title_color=_DT_TEAL)

    # ── SLIDE: Thank You ───────────────────────────────────────────
    end = prs.slides.add_slide(blank)
    ebg = end.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SW, SH)
    _set_shape_fill(ebg, _DT_BLACK)
    ebg.line.fill.background()
    _add_green_bar(end, top=PptxInches(4.0))

    etb = end.shapes.add_textbox(
        PptxInches(0.9), PptxInches(2.5), PptxInches(11), PptxInches(1.2))
    ep = etb.text_frame.paragraphs[0]
    ep.text = "Thank You"
    ep.font.size = PptxPt(44)
    ep.font.bold = True
    ep.font.color.rgb = _DT_WHITE
    ep.alignment = PP_ALIGN.CENTER

    eft = end.shapes.add_textbox(
        PptxInches(0.9), PptxInches(4.3), PptxInches(11), PptxInches(0.6))
    efp = eft.text_frame.paragraphs[0]
    efp.text = "Deloitte AI Advisory  |  Agentic AI ROI Dashboard"
    efp.font.size = PptxPt(14)
    efp.font.color.rgb = _DT_GREEN
    efp.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Render ─────────────────────────────────────────────────────────────────

def _md(text):
    """Escape dollar signs so Streamlit doesn't render them as LaTeX."""
    return str(text).replace("$", r"\$")


def render_roadmap(roadmap):
    """Render a parsed roadmap dict as rich Streamlit components."""
    st.success(_md(roadmap.get("executive_summary", "Roadmap generated successfully.")))

    st.markdown(f"**Total Estimated Cost:** {_md(roadmap.get('total_estimated_cost', 'N/A'))}")
    st.markdown(f"**ROI Timeline:** {_md(roadmap.get('expected_roi_timeline', 'N/A'))}")
    projected = roadmap.get("projected_roi", "")
    if projected:
        st.markdown(f"**Projected ROI:** {_md(projected)}")

    st.markdown("---")

    if roadmap.get("quick_wins"):
        st.subheader("Quick Wins")
        for win in roadmap["quick_wins"]:
            st.markdown(f"- {_md(win)}")

    st.markdown("---")

    for phase in roadmap.get("phases", []):
        pnum = phase.get("phase_number", "?")
        pname = phase.get("phase_name", "Unnamed Phase")
        pdur = phase.get("duration", "")
        pcost = _md(phase.get("estimated_cost", ""))
        with st.expander(
            f"Phase {pnum}: {pname}  ({pdur} · {pcost})",
            expanded=(pnum <= 2 if isinstance(pnum, int) else True),
        ):
            obj_col, del_col = st.columns(2)
            with obj_col:
                st.markdown("**Objectives**")
                for obj in phase.get("objectives", []):
                    st.markdown(f"- {_md(obj)}")
            with del_col:
                deliverables = phase.get("key_deliverables") or phase.get("key_activities", [])
                st.markdown("**Key Deliverables**")
                for d in deliverables:
                    st.markdown(f"- {_md(d)}")

            steps = phase.get("steps", [])
            if steps:
                st.markdown("---")
                st.markdown("**Detailed Implementation Steps**")
                for step in steps:
                    snum = step.get("step_number", "")
                    stitle = _md(step.get("title", ""))
                    sdesc = _md(step.get("description", ""))
                    sdur = step.get("duration", "")
                    sowner = step.get("owner", "")
                    st.markdown(
                        f"**Step {snum}: {stitle}**  \n"
                        f"{sdesc}  \n"
                        f"*Duration: {sdur} · Owner: {sowner}*"
                    )

            res_col, risk_col = st.columns(2)
            with res_col:
                st.markdown("**Resources Needed**")
                for r in phase.get("resources_needed", []):
                    st.markdown(f"- {_md(r)}")
            with risk_col:
                st.markdown("**Risks & Mitigations**")
                for r in phase.get("risks", []):
                    st.markdown(f"- {_md(r)}")

            st.markdown("**Success Metrics**")
            for m in phase.get("success_metrics", []):
                st.markdown(f"- {_md(m)}")

    st.markdown("---")

    if roadmap.get("technology_stack"):
        st.subheader("Recommended Technology Stack")
        cols = st.columns(min(len(roadmap["technology_stack"]), 4))
        for i, tech in enumerate(roadmap["technology_stack"]):
            cols[i % len(cols)].info(tech)

    if roadmap.get("critical_success_factors"):
        st.subheader("Critical Success Factors")
        for f in roadmap["critical_success_factors"]:
            st.markdown(f"- {_md(f)}")

    # ── Assumptions ────────────────────────────────────────────────
    assumptions = roadmap.get("assumptions", [])
    if assumptions:
        st.markdown("---")
        st.subheader("Assumptions")
        st.caption("The following assumptions were made when generating this roadmap. "
                   "Review each and adjust the plan if any assumption does not hold.")
        for a in assumptions:
            if isinstance(a, dict):
                conf = a.get("confidence", "Medium")
                conf_icon = {"High": "🟢", "Medium": "🟡", "Low": "🔴"}.get(conf, "⚪")
                st.markdown(
                    f"- {conf_icon} **{_md(a.get('assumption', ''))}** "
                    f"({conf} confidence)  \n"
                    f"  *If wrong:* {_md(a.get('impact_if_wrong', 'N/A'))}"
                )
            else:
                st.markdown(f"- {_md(a)}")

    # ── Sources ────────────────────────────────────────────────────
    sources = roadmap.get("sources", [])
    if sources:
        st.markdown("---")
        st.subheader("Sources & Attribution")
        st.caption("All data points in this roadmap are attributed to the following sources.")
        for s in sources:
            if isinstance(s, dict):
                ref = _md(s.get("reference", ""))
                url = s.get("url", "")
                used = _md(s.get("used_for", ""))
                if url:
                    st.markdown(f"- **{ref}** — [{url}]({url})  \n  *Used for:* {used}")
                else:
                    st.markdown(f"- **{ref}**  \n  *Used for:* {used}")
            else:
                st.markdown(f"- {_md(s)}")

    # ── Guardrail warnings (output) ────────────────────────────────
    output_warnings = validate_output(roadmap)
    if output_warnings:
        st.markdown("---")
        with st.expander("⚠ Quality Guardrails", expanded=False):
            for w in output_warnings:
                st.warning(w)

    # ── Download buttons ──────────────────────────────────────────────
    st.markdown("---")
    st.subheader("Download Roadmap")

    cache_key = "advisor_downloads"
    if cache_key not in st.session_state or st.session_state.get("_dl_roadmap_id") != id(roadmap):
        exports = {}
        for label, gen_fn, fname, mime in [
            ("PDF", generate_pdf, "ai_implementation_roadmap.pdf", "application/pdf"),
            ("DOCX", generate_docx, "ai_implementation_roadmap.docx",
             "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
            ("PPTX", generate_pptx, "ai_implementation_roadmap.pptx",
             "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        ]:
            try:
                exports[label] = (gen_fn(roadmap), fname, mime)
            except Exception:
                exports[label] = None
        exports["JSON"] = (json.dumps(roadmap, indent=2).encode(),
                           "ai_implementation_roadmap.json", "application/json")
        st.session_state[cache_key] = exports
        st.session_state["_dl_roadmap_id"] = id(roadmap)

    exports = st.session_state[cache_key]
    dl1, dl2, dl3, dl4 = st.columns(4)
    for col, label in [(dl1, "PDF"), (dl2, "DOCX"), (dl3, "PPTX"), (dl4, "JSON")]:
        with col:
            entry = exports.get(label)
            if entry:
                data, fname, mime = entry
                st.download_button(f"Download {label}", data, fname, mime,
                                   use_container_width=True)
